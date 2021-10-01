"""
This file can download slide share presentations on your local system.
without logging in :)
Steps to proceed ->
1). The url of the ppt should be specified in url variables.
2). Using beautiful soup presentation images links are scrapped.
3). Scraped images are downloaded and combined together using pyttx.
4). Presentation is formed.
"""
import io
import os
import requests
from bs4 import BeautifulSoup
from pptx import Presentation

url = 'https://www.slideshare.net/Slideshare/get-started-with-slide-share'

data = requests.get(url)

html = BeautifulSoup(data.text, 'html.parser')
links = html.select('img.slide_image')
prs = Presentation()
for link in links:
    image = link["data-full"]
    response = requests.get(image)
    f = io.BytesIO(response.content)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    placeholder = slide.placeholders[1]
    picture = placeholder.insert_picture(f)

prs.save("Presentation.pptx")
os.startfile("Presentation.pptx")
