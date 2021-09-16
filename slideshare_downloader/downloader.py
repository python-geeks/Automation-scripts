"""
This file can download slide share presentations on your local system.
without logging in :)
Steps to proceed ->
1). The url of the ppt is being provided
2). Using beautiful soup the tag section inside which the images of ppt
    are available  is scraped.
3). Scraped images are combined together using pyttx (forming a presentation)
"""
import io
import os
import requests
from bs4 import BeautifulSoup
from pptx import Presentation

url = 'https://www.slideshare.net/Slideshare/get-started-with-slide-share'

data = requests.get(url)
my_data = []

html = BeautifulSoup(data.text, 'html.parser')
links = html.select('img.slide_image')
prs = Presentation()
for link in links:
    image = link["data-full"]
    my_data.append(image)
    response = requests.get(image)
    f = io.BytesIO(response.content)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    placeholder = slide.placeholders[1]
    picture = placeholder.insert_picture(f)

prs.save("ESEMPIO.pptx")
os.startfile("ESEMPIO.pptx")
print(len(my_data))
